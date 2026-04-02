#!/usr/bin/env python3
"""Generate the Sysadmin Copilot presentation."""

import os
import json
import urllib.request
import urllib.parse
import tempfile

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Image helpers ──────────────────────────────────────────────────────────

_img_cache = {}

def _fetch_image_dalle(prompt, size="1792x1024"):
    api_key = os.environ.get("OPENAI_API_KEY")
    if not api_key:
        return None
    try:
        data = json.dumps({
            "model": "dall-e-3", "prompt": prompt,
            "n": 1, "size": size, "quality": "standard",
        }).encode()
        req = urllib.request.Request(
            "https://api.openai.com/v1/images/generations", data=data,
            headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
        )
        with urllib.request.urlopen(req, timeout=60) as resp:
            result = json.loads(resp.read())
        return _download_image(result["data"][0]["url"])
    except Exception as e:
        print(f"  DALL-E failed: {e}")
        return None

def _fetch_image_pexels(query, orientation="landscape"):
    api_key = os.environ.get("PEXELS_API_KEY")
    if not api_key:
        return None
    try:
        params = urllib.parse.urlencode({"query": query, "orientation": orientation, "size": "large", "per_page": 1})
        req = urllib.request.Request(f"https://api.pexels.com/v1/search?{params}", headers={"Authorization": api_key})
        with urllib.request.urlopen(req, timeout=15) as resp:
            result = json.loads(resp.read())
        if result.get("photos"):
            return _download_image(result["photos"][0]["src"]["large2x"])
    except Exception as e:
        print(f"  Pexels failed: {e}")
    return None

def _fetch_image_unsplash(query):
    access_key = os.environ.get("UNSPLASH_ACCESS_KEY")
    if not access_key:
        return None
    try:
        params = urllib.parse.urlencode({"query": query, "orientation": "landscape", "per_page": 1, "client_id": access_key})
        req = urllib.request.Request(f"https://api.unsplash.com/search/photos?{params}")
        with urllib.request.urlopen(req, timeout=15) as resp:
            result = json.loads(resp.read())
        if result.get("results"):
            return _download_image(result["results"][0]["urls"]["regular"])
    except Exception as e:
        print(f"  Unsplash failed: {e}")
    return None

def _download_image(url):
    if url in _img_cache:
        return _img_cache[url]
    tmp = tempfile.NamedTemporaryFile(suffix=".jpg", delete=False)
    urllib.request.urlretrieve(url, tmp.name)
    _img_cache[url] = tmp.name
    return tmp.name

def fetch_image(prompt, search_query=None):
    if search_query is None:
        search_query = prompt
    if prompt in _img_cache:
        return _img_cache[prompt]
    print(f"  Fetching image: {search_query}...")
    path = _fetch_image_dalle(prompt) or _fetch_image_pexels(search_query) or _fetch_image_unsplash(search_query)
    if path:
        _img_cache[prompt] = path
    return path

def add_slide_image(slide, img_path, left, top, width, height=None):
    if height:
        slide.shapes.add_picture(img_path, Inches(left), Inches(top), Inches(width), Inches(height))
    else:
        slide.shapes.add_picture(img_path, Inches(left), Inches(top), Inches(width))

def add_fullbleed_image(slide, img_path):
    pic = slide.shapes.add_picture(img_path, 0, 0, prs.slide_width, prs.slide_height)
    sp = pic._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return pic

def cleanup_images():
    for path in _img_cache.values():
        try:
            os.unlink(path)
        except OSError:
            pass

# ── Core helpers ───────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette — light vivid
BG_DARK    = RGBColor(0xF5, 0xF7, 0xFA)  # soft warm white
BG_CARD    = RGBColor(0xE8, 0xEC, 0xF2)  # light card bg
ACCENT     = RGBColor(0x26, 0x6E, 0xF1)  # vivid blue
ACCENT2    = RGBColor(0x0F, 0xB5, 0x5E)  # vivid green
ACCENT3    = RGBColor(0xF5, 0x6E, 0x0F)  # vivid orange
RED        = RGBColor(0xE8, 0x31, 0x31)  # vivid red
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)  # pure white (for cards/contrast)
LIGHT_GRAY = RGBColor(0x3A, 0x3F, 0x4C)  # body text (dark on light bg)
MID_GRAY   = RGBColor(0x5A, 0x5F, 0x6C)  # secondary text
DIM_GRAY   = RGBColor(0x7A, 0x80, 0x8C)  # muted text
CODE_BG    = RGBColor(0x1E, 0x29, 0x3B)  # dark code boxes for contrast
TITLE_TEXT = RGBColor(0x1A, 0x1D, 0x2E)  # near-black for titles

def add_slide(layout_index=6):
    return prs.slides.add_slide(prs.slide_layouts[layout_index])

def add_shape_bg(slide, color_rgb):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color_rgb
    shape.line.fill.background()
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return shape

def add_textbox(slide, left, top, width, height, text,
                font_size=18, bold=False, color=LIGHT_GRAY,
                alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_rich_textbox(slide, left, top, width, height):
    """Return a text_frame for multi-paragraph use."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf

def add_para(tf, text, font_size=18, bold=False, color=LIGHT_GRAY,
             font_name="Calibri", alignment=PP_ALIGN.LEFT, space_before=0, space_after=4):
    """Add a paragraph. Default color is body text (dark on light bg)."""
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    return p

def add_code_box(slide, left, top, width, height, text):
    """Dark rounded rectangle with monospace text."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.2)
    tf.margin_right = Inches(0.3)
    tf.margin_bottom = Inches(0.2)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(0xA0, 0xE8, 0x70)  # green code
    p.font.name = "Consolas"
    p.alignment = PP_ALIGN.LEFT
    return shape

def add_accent_line(slide, top, color=ACCENT):
    """Thin accent line across the slide."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(top), Inches(11.7), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, bg=BG_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_circle_icon(slide, left, top, size, color, text):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(int(size * 18))
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    return shape

def slide_number_label(slide, num):
    add_textbox(slide, 12.2, 7.0, 1, 0.4, str(num), font_size=11, color=DIM_GRAY, alignment=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════
print("Building slides...")

s = add_slide()
add_shape_bg(s, BG_DARK)
# Try hero image
img = fetch_image(
    prompt="Dark futuristic terminal console with blue glowing code on a black screen, minimalist, cinematic lighting, no text",
    search_query="terminal code dark"
)
if img:
    add_fullbleed_image(s, img)
    # dark overlay
    overlay = add_shape_bg(s, RGBColor(0xF5, 0xF7, 0xFA))

# Accent bar at top
shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT; shape.line.fill.background()

add_textbox(s, 1, 1.8, 11.3, 1.2, "SYSADMIN COPILOT", font_size=56, bold=True, color=TITLE_TEXT, alignment=PP_ALIGN.CENTER, font_name="Calibri")
add_textbox(s, 1, 3.1, 11.3, 0.8, "Talk to Your Linux Server in Plain English", font_size=28, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_accent_line(s, 4.1)
add_textbox(s, 1, 4.5, 11.3, 0.6, "An AI Agent Built with LangChain + LangGraph", font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Terminal prompt icon
add_code_box(s, 4.5, 5.5, 4.3, 0.6, "$ sysadmin-copilot")

print("  Slide 1: Title")

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — AGENDA
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); add_shape_bg(s, BG_DARK)
add_textbox(s, 0.8, 0.4, 11, 0.8, "What We'll Cover Today", font_size=38, bold=True, color=TITLE_TEXT)
add_accent_line(s, 1.2)

agenda = [
    ("The Problem", "Why sysadmin work is hard", "5 min", ACCENT),
    ("The Solution", "An AI agent that runs CLI tools for you", "5 min", ACCENT),
    ("Architecture", "How the pieces fit together", "10 min", ACCENT2),
    ("Deep Dive", "Agent, Tools, Safety, Audit", "15 min", ACCENT2),
    ("Live Demo", "See it in action", "10 min", ACCENT3),
    ("Workshop", "Build your own plugin tool", "15 min", ACCENT3),
]
y = 1.7
for title, desc, time, clr in agenda:
    add_card(s, 0.8, y, 9.5, 0.65)
    add_textbox(s, 1.1, y + 0.05, 2.5, 0.5, title, font_size=20, bold=True, color=clr)
    add_textbox(s, 3.6, y + 0.05, 5.5, 0.5, desc, font_size=18, color=LIGHT_GRAY)
    add_textbox(s, 10.5, y + 0.05, 1.8, 0.5, time, font_size=16, color=DIM_GRAY, alignment=PP_ALIGN.RIGHT)
    y += 0.82

add_textbox(s, 0.8, y + 0.3, 11, 0.5, "Total: ~60 minutes", font_size=16, color=DIM_GRAY, alignment=PP_ALIGN.RIGHT)
slide_number_label(s, 2)
print("  Slide 2: Agenda")

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); add_shape_bg(s, BG_DARK)
add_textbox(s, 0.8, 0.4, 11, 0.8, "The Sysadmin's Daily Struggle", font_size=38, bold=True, color=TITLE_TEXT)
add_accent_line(s, 1.2)

# Quote-style pain points
pains = [
    '"What journalctl flags do I need again?"',
    '"Was it ss or netstat? What\'s the syntax?"',
    '"The server is slow \u2014 where do I even start?"',
]
y = 1.8
for pain in pains:
    add_card(s, 1.0, y, 6.5, 0.65)
    add_textbox(s, 1.3, y + 0.08, 6, 0.5, pain, font_size=19, color=ACCENT, font_name="Calibri")
    y += 0.85

tf = add_rich_textbox(s, 1.0, y + 0.3, 6.5, 1.5)
add_para(tf, "Knowledge takes years to accumulate", font_size=20, color=LIGHT_GRAY, space_after=8)
add_para(tf, "Even experts spend time looking up flags and syntax", font_size=20, color=LIGHT_GRAY)

# Right side image
img = fetch_image(
    prompt="A frustrated system administrator staring at a terminal with error messages, dark moody lighting, no text",
    search_query="programmer frustrated computer"
)
if img:
    add_slide_image(s, img, 8.2, 1.5, 4.5, 5.0)
else:
    add_card(s, 8.2, 1.5, 4.5, 5.0)
    add_textbox(s, 8.7, 3.5, 3.5, 1, "?!", font_size=72, bold=True, color=DIM_GRAY, alignment=PP_ALIGN.CENTER)

slide_number_label(s, 3)
print("  Slide 3: Problem")

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — THE SOLUTION
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); add_shape_bg(s, BG_DARK)
add_textbox(s, 0.8, 0.4, 11, 0.8, "What If You Could Just Ask?", font_size=38, bold=True, color=TITLE_TEXT)
add_accent_line(s, 1.2)

code_text = (
    "\u276f Why is the server running slow?\n"
    "\n"
    "  [check_cpu_and_load]\n"
    "  [check_top_processes]\n"
    "  [check_memory]\n"
    "\n"
    "The server load average is 4.2 \u2014 high for 2 cores.\n"
    "Top CPU consumer: python3 at 87%.\n"
    "Memory is fine: 3.1 GB used of 8 GB."
)
add_code_box(s, 1.0, 1.6, 7.5, 3.8, code_text)

add_card(s, 9.0, 1.6, 3.8, 3.8)
tf = add_rich_textbox(s, 9.3, 1.9, 3.2, 3.2)
add_para(tf, "The agent:", font_size=20, bold=True, color=ACCENT, space_after=12)
add_para(tf, "1. Reasons about your question", font_size=17, color=LIGHT_GRAY, space_after=8)
add_para(tf, "2. Picks the right tools", font_size=17, color=LIGHT_GRAY, space_after=8)
add_para(tf, "3. Runs them in sequence", font_size=17, color=LIGHT_GRAY, space_after=8)
add_para(tf, "4. Explains the results", font_size=17, color=LIGHT_GRAY)

add_textbox(s, 1.0, 5.8, 11.5, 0.6, "You describe the problem. The AI figures out the commands.", font_size=22, bold=True, color=ACCENT2, alignment=PP_ALIGN.CENTER)
slide_number_label(s, 4)
print("  Slide 4: Solution")

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — WHO IT'S FOR
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); add_shape_bg(s, BG_DARK)
add_textbox(s, 0.8, 0.4, 11, 0.8, "Who Is This For?", font_size=38, bold=True, color=TITLE_TEXT)
add_accent_line(s, 1.2)

cols = [
    ("Homelab &\nSelf-Hosters", "Manage servers without being a full-time sysadmin", ACCENT, 0.8),
    ("Developers", "Investigate production issues without Linux tool expertise", ACCENT2, 4.8),
    ("Sysadmins", "Faster first-pass investigation with built-in audit trail", ACCENT3, 8.8),
]
for title, desc, clr, x in cols:
    add_card(s, x, 1.7, 3.7, 3.5)
    add_circle_icon(s, x + 1.35, 2.0, 1.0, clr, title[0])
    add_textbox(s, x + 0.2, 3.2, 3.3, 0.8, title, font_size=20, bold=True, color=clr, alignment=PP_ALIGN.CENTER)
    add_textbox(s, x + 0.2, 4.0, 3.3, 1.0, desc, font_size=15, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(s, 0.8, 5.7, 11.5, 0.5, "Targets Linux with systemd  (journalctl, systemctl, ss, dig, free ...)", font_size=16, color=DIM_GRAY, alignment=PP_ALIGN.CENTER)
slide_number_label(s, 5)
print("  Slide 5: Who it's for")

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — KEY FEATURES
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); add_shape_bg(s, BG_DARK)
add_textbox(s, 0.8, 0.4, 11, 0.8, "Feature Highlights", font_size=38, bold=True, color=TITLE_TEXT)
add_accent_line(s, 1.2)

features = [
    ("NL", "Natural Language\nInterface", ACCENT),
    ("RA", "ReAct Agent Loop\n(Reason + Act)", ACCENT),
    ("MT", "Multi-turn\nConversation", ACCENT2),
    ("RWB", "Three-tier Safety\nREAD / WRITE / BLOCKED", ACCENT3),
    ("LOG", "Full Audit Trail\n(JSONL)", ACCENT2),
    ("LLM", "Multiple Backends\nOllama / OpenAI / Anthropic", ACCENT),
    ("PLG", "Plugin System\ntools_extra/", ACCENT3),
]

cols_per_row = 4
x_start = 0.6
y_start = 1.6
card_w = 2.9
card_h = 2.2
gap = 0.25

for i, (icon, label, clr) in enumerate(features):
    row = i // cols_per_row
    col = i % cols_per_row
    x = x_start + col * (card_w + gap)
    y = y_start + row * (card_h + 0.3)
    add_card(s, x, y, card_w, card_h)
    add_circle_icon(s, x + (card_w - 0.7) / 2, y + 0.25, 0.7, clr, icon)
    add_textbox(s, x + 0.15, y + 1.15, card_w - 0.3, 0.9, label, font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

slide_number_label(s, 6)
print("  Slide 6: Features")

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); add_shape_bg(s, BG_DARK)
add_textbox(s, 0.8, 0.4, 11, 0.8, "Architecture: The Big Picture", font_size=38, bold=True, color=TITLE_TEXT)
add_accent_line(s, 1.2)

# Flow boxes: agent.py -> LangGraph -> safety.py -> tools.py / audit.py
boxes = [
    (0.5, 2.5, 2.4, 1.8, "agent.py", "Entry Point\nREPL Loop", ACCENT),
    (3.4, 2.5, 2.4, 1.8, "LangGraph", "Reason \u2192 Act\n\u2192 Reason", ACCENT2),
    (6.3, 2.5, 2.4, 1.8, "safety.py", "Permission\nGates", ACCENT3),
    (9.2, 2.0, 2.4, 1.3, "tools.py", "27 CLI Tools", ACCENT),
    (9.2, 3.6, 2.4, 1.3, "audit.py", "Every Call\nLogged", ACCENT2),
]

for x, y, w, h, title, desc, clr in boxes:
    add_card(s, x, y, w, h)
    add_textbox(s, x + 0.1, y + 0.15, w - 0.2, 0.5, title, font_size=18, bold=True, color=clr, alignment=PP_ALIGN.CENTER)
    add_textbox(s, x + 0.1, y + 0.65, w - 0.2, h - 0.8, desc, font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Arrows between boxes
for ax in [2.9, 5.8, 8.7]:
    arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(ax), Inches(3.15), Inches(0.5), Inches(0.3))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = DIM_GRAY; arrow.line.fill.background()

# Fork arrow from safety to tools/audit
arrow2 = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.7), Inches(3.85), Inches(0.5), Inches(0.3))
arrow2.fill.solid(); arrow2.fill.fore_color.rgb = DIM_GRAY; arrow2.line.fill.background()

add_textbox(s, 0.5, 5.2, 11, 0.5, "Four modules, four responsibilities. Each does one thing well.", font_size=18, color=DIM_GRAY, alignment=PP_ALIGN.CENTER)
slide_number_label(s, 7)
print("  Slide 7: Architecture")

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — DATA FLOW
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); add_shape_bg(s, BG_DARK)
add_textbox(s, 0.8, 0.4, 11, 0.8, "What Happens When You Ask a Question", font_size=36, bold=True, color=TITLE_TEXT)
add_accent_line(s, 1.2)

steps = [
    ("1", 'User types: "Check if nginx is running"', ACCENT),
    ("2", "Message appended to conversation history", LIGHT_GRAY),
    ("3", 'LLM reasons: "I should call check_service_status"', ACCENT2),
    ("4", "Safety layer checks for blocked patterns", ACCENT3),
    ("5", "Tool runs: systemctl status nginx", ACCENT),
    ("6", "Output returned to LLM (capped at 8000 chars)", LIGHT_GRAY),
    ("7", "LLM generates a plain-English answer", ACCENT2),
    ("8", "Answer streamed token-by-token to terminal", ACCENT),
]

y = 1.5
for num, text, clr in steps:
    add_circle_icon(s, 1.0, y, 0.5, clr, num)
    add_textbox(s, 1.7, y + 0.02, 10, 0.5, text, font_size=18, color=LIGHT_GRAY)
    y += 0.68

slide_number_label(s, 8)
print("  Slide 8: Data flow")

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — THE AGENT
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); add_shape_bg(s, BG_DARK)
add_textbox(s, 0.8, 0.4, 11, 0.8, "The Agent: ReAct Pattern", font_size=38, bold=True, color=TITLE_TEXT)
add_accent_line(s, 1.2)

# Left side explanation
tf = add_rich_textbox(s, 0.8, 1.6, 5.5, 4.5)
add_para(tf, "ReAct = Reason + Act", font_size=22, bold=True, color=ACCENT, space_after=14)
add_para(tf, "Thought \u2192 Action \u2192 Observation \u2192 Thought", font_size=18, color=ACCENT2, space_after=14)
add_para(tf, "Built on LangChain's create_agent + LangGraph", font_size=17, color=LIGHT_GRAY, space_after=8)
add_para(tf, "System prompt includes hostname, OS, time", font_size=17, color=LIGHT_GRAY, space_after=8)
add_para(tf, "Anti-hallucination directives for weaker models", font_size=17, color=LIGHT_GRAY, space_after=8)
add_para(tf, '"Only report info from tool output"', font_size=17, color=DIM_GRAY, space_after=8)
add_para(tf, '"Prefer specific tools over run_command"', font_size=17, color=DIM_GRAY)

# Right side: code
code = (
    "llm = get_llm()\n"
    "\n"
    "wrapped_tools = safety.wrap_tools(\n"
    "    ALL_TOOLS, audit\n"
    ")\n"
    "\n"
    "agent = create_agent(\n"
    "    model=llm,\n"
    "    tools=wrapped_tools,\n"
    "    system_prompt=build_system_prompt(),\n"
    ")"
)
add_code_box(s, 6.8, 1.6, 5.8, 3.5, code)
add_textbox(s, 6.8, 5.2, 5.8, 0.4, "agent.py \u2014 three steps to launch", font_size=13, color=DIM_GRAY, alignment=PP_ALIGN.CENTER)

slide_number_label(s, 9)
print("  Slide 9: Agent")

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — CONVERSATION HISTORY
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); add_shape_bg(s, BG_DARK)
add_textbox(s, 0.8, 0.4, 11, 0.8, "Multi-Turn Memory", font_size=38, bold=True, color=TITLE_TEXT)
add_accent_line(s, 1.2)

conv = (
    "\u276f Show failed SSH logins in the last hour\n"
    "  [query_journal_logs]\n"
    "  Found 23 failed attempts from 185.220.101.47\n"
    "\n"
    "\u276f Is that IP doing anything else?\n"
    "  [check_network_connections]\n"
    "  No active connections from that IP\n"
    "\n"
    "\u276f Restart fail2ban\n"
    "  \u26a0  Allow this action? [y/N]: y\n"
    "  [restart_service]\n"
    "  fail2ban restarted and active"
)
add_code_box(s, 1.0, 1.6, 7.0, 4.5, conv)

add_card(s, 8.5, 1.6, 4.3, 4.5)
tf = add_rich_textbox(s, 8.8, 1.9, 3.7, 4.0)
add_para(tf, "Key insight:", font_size=20, bold=True, color=ACCENT, space_after=12)
add_para(tf, '"that IP" just works \u2014 context carries across turns', font_size=17, color=LIGHT_GRAY, space_after=16)
add_para(tf, "Three different tools were called automatically", font_size=17, color=LIGHT_GRAY, space_after=16)
add_para(tf, "Service restart required explicit y/N confirmation", font_size=17, color=ACCENT3, space_after=16)
add_para(tf, 'Type "new" to reset history', font_size=15, color=DIM_GRAY)

slide_number_label(s, 10)
print("  Slide 10: Conversation")

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 11 — TOOLS
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); add_shape_bg(s, BG_DARK)
add_textbox(s, 0.8, 0.4, 11, 0.8, "27 Tools in 7 Categories", font_size=38, bold=True, color=TITLE_TEXT)
add_accent_line(s, 1.2)

categories = [
    ("Logs", "journal, log files, dmesg", ACCENT, 0.8, 1.6),
    ("System Health", "disk, memory, CPU, processes, zombies", ACCENT2, 0.8, 2.5),
    ("Services", "status, failed, restart, stop", ACCENT3, 0.8, 3.4),
    ("Network", "ports, connections, ping, DNS, URL health", ACCENT, 0.8, 4.3),
    ("Users & Files", "logged-in users, cron jobs, recent files", ACCENT2, 6.8, 1.6),
    ("Security", "system audit, outdated packages, updates", ACCENT3, 6.8, 2.5),
    ("General", "run_command, change_directory, web search", DIM_GRAY, 6.8, 3.4),
]

for name, tools, clr, x, y in categories:
    add_card(s, x, y, 5.7, 0.75)
    add_textbox(s, x + 0.15, y + 0.05, 2.0, 0.4, name, font_size=17, bold=True, color=clr)
    add_textbox(s, x + 2.2, y + 0.05, 3.3, 0.6, tools, font_size=14, color=LIGHT_GRAY)

add_textbox(s, 0.8, 5.4, 11.5, 0.6, "Each tool is a @tool-decorated function with a docstring the LLM reads to decide when to use it.", font_size=16, color=DIM_GRAY, alignment=PP_ALIGN.CENTER)
slide_number_label(s, 11)
print("  Slide 11: Tools")

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 12 — run_cmd()
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); add_shape_bg(s, BG_DARK)
add_textbox(s, 0.8, 0.4, 11, 0.8, "run_cmd(): The Engine Room", font_size=38, bold=True, color=TITLE_TEXT)
add_accent_line(s, 1.2)

decisions = [
    ("List-based subprocess.run()", "No shell injection \u2014 args go directly to the kernel", ACCENT2),
    ("capture_output=True", "Captures both stdout and stderr", LIGHT_GRAY),
    ("30-second timeout", "No hanging on slow commands like ping", ACCENT3),
    ("8000 char output cap", "Protects the LLM's context window", ACCENT),
    ("Graceful error handling", "Timeouts, missing commands, permission errors", LIGHT_GRAY),
]

y = 1.6
for title, desc, clr in decisions:
    add_card(s, 0.8, y, 6.2, 0.7)
    add_textbox(s, 1.1, y + 0.05, 6, 0.3, title, font_size=17, bold=True, color=clr)
    add_textbox(s, 1.1, y + 0.35, 6, 0.3, desc, font_size=14, color=DIM_GRAY)
    y += 0.85

code = (
    "def run_cmd(cmd: list[str],\n"
    "            timeout: int = 30) -> str:\n"
    "    result = subprocess.run(\n"
    "        cmd,\n"
    "        capture_output=True,\n"
    "        text=True,\n"
    "        timeout=timeout,\n"
    "    )\n"
    "    # truncate if > MAX_OUTPUT_CHARS\n"
    "    # append stderr on failure\n"
    "    return output"
)
add_code_box(s, 7.5, 1.6, 5.2, 3.5, code)

slide_number_label(s, 12)
print("  Slide 12: run_cmd")

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 13 — SAFETY LAYER
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); add_shape_bg(s, BG_DARK)
add_textbox(s, 0.8, 0.4, 11, 0.8, "Three-Tier Safety Model", font_size=38, bold=True, color=TITLE_TEXT)
add_accent_line(s, 1.2)

tiers = [
    ("READ", "Always allowed", "Log queries, status checks, disk info", ACCENT2, 1.6),
    ("WRITE", "Confirmation required", "restart_service, stop_service, update_packages", ACCENT3, 3.0),
    ("BLOCKED", "Rejected immediately", "rm, dd, shutdown, reboot, fork bombs", RED, 4.4),
]

for name, subtitle, desc, clr, y in tiers:
    # Colored bar on left
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(y), Inches(0.15), Inches(1.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = clr; bar.line.fill.background()
    add_card(s, 1.1, y, 11.4, 1.1)
    add_textbox(s, 1.4, y + 0.05, 2.5, 0.5, name, font_size=26, bold=True, color=clr)
    add_textbox(s, 1.4, y + 0.55, 2.5, 0.4, subtitle, font_size=15, color=LIGHT_GRAY)
    add_textbox(s, 4.2, y + 0.25, 8, 0.6, desc, font_size=17, color=DIM_GRAY)

add_textbox(s, 0.8, 5.9, 11.5, 0.5, "Rules run in Python, not in the prompt \u2014 the LLM cannot override them.", font_size=19, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
slide_number_label(s, 13)
print("  Slide 13: Safety")

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 14 — BLOCKED PATTERNS
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); add_shape_bg(s, BG_DARK)
add_textbox(s, 0.8, 0.4, 11, 0.8, "Defense in Depth: Blocked Patterns", font_size=36, bold=True, color=TITLE_TEXT)
add_accent_line(s, 1.2)

patterns = [
    ("File destruction", "rm, rmdir, shred, truncate, unlink"),
    ("Disk / device", 'dd, mkfs, > /dev/, tee /dev/'),
    ("System state", "shutdown, reboot, poweroff, halt"),
    ("Permissions", "chmod 777, chmod a+rwx"),
    ("Evasion", "base64 decode, | bash, | sh"),
]

y = 1.6
for cat, pats in patterns:
    add_card(s, 0.8, y, 7, 0.6)
    add_textbox(s, 1.1, y + 0.08, 2.2, 0.4, cat, font_size=16, bold=True, color=RED)
    add_textbox(s, 3.3, y + 0.08, 4.3, 0.4, pats, font_size=14, color=LIGHT_GRAY, font_name="Consolas")
    y += 0.72

add_card(s, 8.3, 1.6, 4.5, 3.5)
tf = add_rich_textbox(s, 8.6, 1.9, 3.9, 3.0)
add_para(tf, "Normalization", font_size=19, bold=True, color=ACCENT3, space_after=12)
add_para(tf, "lowercase  \u2192  catches RM, Rm", font_size=15, color=LIGHT_GRAY, space_after=6, font_name="Consolas")
add_para(tf, "collapse whitespace  \u2192  rm\\tfile", font_size=15, color=LIGHT_GRAY, space_after=6, font_name="Consolas")
add_para(tf, 'strip quotes  \u2192  "rm" file', font_size=15, color=LIGHT_GRAY, space_after=12, font_name="Consolas")
add_para(tf, "Defense-in-depth, not a security boundary.", font_size=14, color=DIM_GRAY, space_after=4)
add_para(tf, "OS permissions are the real boundary.", font_size=14, bold=True, color=ACCENT)

slide_number_label(s, 14)
print("  Slide 14: Blocked patterns")

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 15 — THREE LAYER DEFENSE
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); add_shape_bg(s, BG_DARK)
add_textbox(s, 0.8, 0.4, 11, 0.8, "The Real Security Model", font_size=38, bold=True, color=TITLE_TEXT)
add_accent_line(s, 1.2)

layers = [
    ("Layer 1", "Blocked Patterns", "Catches common mistakes\nand obvious attacks", ACCENT3, 1.7),
    ("Layer 2", "User Confirmation", "Human approves every\nWRITE action before execution", ACCENT, 3.0),
    ("Layer 3", "OS Permissions", "Service account + sudoers\n= the real security boundary", ACCENT2, 4.3),
]

for label, name, desc, clr, y in layers:
    # Full-width card
    add_card(s, 0.8, y, 11.7, 1.1)
    add_circle_icon(s, 1.2, y + 0.17, 0.75, clr, label[-1])
    add_textbox(s, 2.2, y + 0.1, 3.5, 0.5, name, font_size=22, bold=True, color=clr)
    add_textbox(s, 2.2, y + 0.55, 3.5, 0.5, label, font_size=14, color=DIM_GRAY)
    add_textbox(s, 6.5, y + 0.15, 5.5, 0.8, desc, font_size=17, color=LIGHT_GRAY)

add_textbox(s, 0.8, 5.8, 11.5, 0.5, "Even if layers 1 and 2 fail, the OS won't allow it.", font_size=20, bold=True, color=ACCENT2, alignment=PP_ALIGN.CENTER)
slide_number_label(s, 15)
print("  Slide 15: Three layers")

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 16 — AUDIT LOGGER
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); add_shape_bg(s, BG_DARK)
add_textbox(s, 0.8, 0.4, 11, 0.8, "Full Audit Trail", font_size=38, bold=True, color=TITLE_TEXT)
add_accent_line(s, 1.2)

jsonl = (
    '{"tool": "check_service_status",\n'
    ' "args": {"service": "nginx"}, "status": "OK"}\n'
    '\n'
    '{"tool": "restart_service",\n'
    ' "args": {"service": "nginx"}, "status": "CONFIRMED"}\n'
    '\n'
    '{"tool": "stop_service",\n'
    ' "args": {"service": "mysql"}, "status": "DENIED"}'
)
add_code_box(s, 0.8, 1.6, 7.5, 3.2, jsonl)

# Status legend
add_card(s, 8.8, 1.6, 4.0, 3.2)
statuses = [
    ("OK", "Read-only tool call", ACCENT2),
    ("CONFIRMED", "User approved write", ACCENT),
    ("DENIED", "User rejected write", ACCENT3),
    ("BLOCKED", "Caught by safety layer", RED),
]
y = 1.9
for status, desc, clr in statuses:
    add_textbox(s, 9.1, y, 1.7, 0.4, status, font_size=15, bold=True, color=clr, font_name="Consolas")
    add_textbox(s, 10.5, y, 2.1, 0.4, desc, font_size=14, color=LIGHT_GRAY)
    y += 0.6

# REPL commands
add_card(s, 0.8, 5.2, 11.7, 1.0)
tf = add_rich_textbox(s, 1.1, 5.35, 11, 0.7)
add_para(tf, "audit              \u2192  current session      |      audit last          \u2192  previous session      |      audit last 3   \u2192  last 3 sessions", font_size=15, color=LIGHT_GRAY, font_name="Consolas")

slide_number_label(s, 16)
print("  Slide 16: Audit")

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 17 — CONFIGURATION
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); add_shape_bg(s, BG_DARK)
add_textbox(s, 0.8, 0.4, 11, 0.8, "Configuration: One Env Var to Switch", font_size=36, bold=True, color=TITLE_TEXT)
add_accent_line(s, 1.2)

backends = (
    "python agent.py                       # Ollama (local, default)\n"
    "LLM_PROVIDER=openai python agent.py   # OpenAI\n"
    "LLM_PROVIDER=anthropic python agent.py # Anthropic"
)
add_code_box(s, 0.8, 1.6, 11.7, 1.5, backends)

# Env var table
env_vars = [
    ("LLM_PROVIDER", "ollama", "Backend selection"),
    ("OLLAMA_MODEL", "qwen3:8b", "Ollama model name"),
    ("OPENAI_MODEL", "gpt-4o-mini", "OpenAI model"),
    ("ANTHROPIC_MODEL", "claude-sonnet-4", "Anthropic model"),
    ("EXTRA_SERVICES", "\u2014", "Add services to allowlist"),
    ("LOG_PATHS", "/var/log", "Allowed log directories"),
]

# Table header
add_card(s, 0.8, 3.5, 11.7, 0.5, bg=RGBColor(0xD8, 0xDE, 0xE9))
add_textbox(s, 1.0, 3.55, 3.5, 0.4, "Variable", font_size=15, bold=True, color=ACCENT, font_name="Consolas")
add_textbox(s, 4.5, 3.55, 3, 0.4, "Default", font_size=15, bold=True, color=ACCENT, font_name="Consolas")
add_textbox(s, 7.5, 3.55, 4.5, 0.4, "Purpose", font_size=15, bold=True, color=ACCENT)

y = 4.1
for var, default, purpose in env_vars:
    add_textbox(s, 1.0, y, 3.5, 0.35, var, font_size=14, color=LIGHT_GRAY, font_name="Consolas")
    add_textbox(s, 4.5, y, 3, 0.35, default, font_size=14, color=DIM_GRAY, font_name="Consolas")
    add_textbox(s, 7.5, y, 4.5, 0.35, purpose, font_size=14, color=LIGHT_GRAY)
    y += 0.42

slide_number_label(s, 17)
print("  Slide 17: Configuration")

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 18 — INSTALLATION
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); add_shape_bg(s, BG_DARK)
add_textbox(s, 0.8, 0.4, 11, 0.8, "One Command Install", font_size=38, bold=True, color=TITLE_TEXT)
add_accent_line(s, 1.2)

add_code_box(s, 3.5, 1.5, 6, 0.65, "$ sudo bash install.sh")

steps = [
    ("1", "Creates sysadmin-copilot service account", ACCENT),
    ("2", "Sets up groups (systemd-journal, adm)", ACCENT2),
    ("3", "Installs app to /opt/sysadmin-copilot/", ACCENT),
    ("4", "Configures sudoers (validated with visudo!)", ACCENT3),
    ("5", "Creates /usr/local/bin/sysadmin-copilot wrapper", ACCENT2),
]

y = 2.6
for num, text, clr in steps:
    add_circle_icon(s, 1.5, y, 0.5, clr, num)
    add_textbox(s, 2.2, y + 0.03, 9, 0.5, text, font_size=19, color=LIGHT_GRAY)
    y += 0.7

add_code_box(s, 3.5, y + 0.3, 6, 0.65, "$ sysadmin-copilot    # that's it!")
add_textbox(s, 0.8, y + 1.2, 11.5, 0.4, "Least-privilege service account \u2014 contained blast radius even if compromised", font_size=16, color=DIM_GRAY, alignment=PP_ALIGN.CENTER)

slide_number_label(s, 18)
print("  Slide 18: Installation")

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 19 — PLUGIN SYSTEM
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); add_shape_bg(s, BG_DARK)
add_textbox(s, 0.8, 0.4, 11, 0.8, "Extending: The Plugin System", font_size=38, bold=True, color=TITLE_TEXT)
add_accent_line(s, 1.2)

tf = add_rich_textbox(s, 0.8, 1.6, 6, 4)
add_para(tf, "How it works:", font_size=22, bold=True, color=ACCENT, space_after=14)
add_para(tf, "Drop a .py file in tools_extra/", font_size=18, color=LIGHT_GRAY, space_after=8)
add_para(tf, "Auto-discovered at startup", font_size=18, color=LIGHT_GRAY, space_after=8)
add_para(tf, "No core files need editing", font_size=18, color=ACCENT2, bold=True, space_after=8)
add_para(tf, "Declare WRITE_TOOLS for write actions", font_size=18, color=LIGHT_GRAY, space_after=8)
add_para(tf, "Same safety wrapping as core tools", font_size=18, color=LIGHT_GRAY, space_after=8)
add_para(tf, "Errors warned, don't crash startup", font_size=18, color=LIGHT_GRAY)

plugin_code = (
    '# tools_extra/docker_tools.py\n'
    'from langchain_core.tools import tool\n'
    'from tools import run_cmd\n'
    '\n'
    '@tool\n'
    'def check_docker_containers(\n'
    '    all: bool = False\n'
    ') -> str:\n'
    '    """List Docker containers."""\n'
    '    cmd = ["docker", "ps"]\n'
    '    if all:\n'
    '        cmd.append("-a")\n'
    '    return run_cmd(cmd)\n'
    '\n'
    'WRITE_TOOLS = set()'
)
add_code_box(s, 7.2, 1.5, 5.6, 4.5, plugin_code)

add_code_box(s, 0.8, 5.5, 11.7, 0.55, "Loaded 2 extra tool(s): check_docker_containers, check_docker_images")

slide_number_label(s, 19)
print("  Slide 19: Plugins")

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 20 — WORKSHOP DIVIDER
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); add_shape_bg(s, BG_DARK)

img = fetch_image(
    prompt="Dark abstract cyber security visualization with network nodes and glowing connections, no text, cinematic",
    search_query="cyber security network abstract"
)
if img:
    add_fullbleed_image(s, img)
    add_shape_bg(s, RGBColor(0xF5, 0xF7, 0xFA))

# Top accent bar
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT3; bar.line.fill.background()

add_textbox(s, 1, 2.2, 11.3, 1.2, "WORKSHOP TIME!", font_size=60, bold=True, color=ACCENT3, alignment=PP_ALIGN.CENTER)
add_textbox(s, 1, 3.6, 11.3, 0.8, "Let's Build a Threat Intelligence Plugin", font_size=30, color=TITLE_TEXT, alignment=PP_ALIGN.CENTER)

# Decorative line
add_accent_line(s, 4.6, ACCENT3)

add_textbox(s, 1, 5.0, 11.3, 0.6, "tools_extra/threat_intel.py", font_size=22, color=ACCENT, alignment=PP_ALIGN.CENTER, font_name="Consolas")

print("  Slide 20: Workshop divider")

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 21 — WORKSHOP: WHAT WE'RE BUILDING
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); add_shape_bg(s, BG_DARK)
add_textbox(s, 0.8, 0.4, 11, 0.8, "Workshop: Threat Intel Plugin", font_size=36, bold=True, color=ACCENT3)
add_accent_line(s, 1.2, ACCENT3)

tools_list = [
    ("extract_iocs", "Extract IPs, domains, URLs, hashes from log files", ACCENT),
    ("hash_file", "Compute MD5 / SHA1 / SHA256 of any file", ACCENT2),
    ("vt_hash_lookup", "Look up file hash on VirusTotal", ACCENT3),
    ("vt_ip_lookup", "Look up IP reputation on VirusTotal", ACCENT3),
    ("vt_domain_lookup", "Look up domain reputation on VirusTotal", ACCENT3),
]

y = 1.6
for name, desc, clr in tools_list:
    add_card(s, 0.8, y, 11.7, 0.7)
    add_textbox(s, 1.1, y + 0.1, 3.2, 0.4, name, font_size=18, bold=True, color=clr, font_name="Consolas")
    add_textbox(s, 4.5, y + 0.1, 7.5, 0.5, desc, font_size=17, color=LIGHT_GRAY)
    y += 0.85

add_card(s, 3.5, y + 0.2, 6.3, 0.6, bg=RGBColor(0xE0, 0xF5, 0xE8))
add_textbox(s, 3.7, y + 0.3, 5.9, 0.4, "All read-only tools \u2014 WRITE_TOOLS = {} (empty)", font_size=17, bold=True, color=ACCENT2, alignment=PP_ALIGN.CENTER)

slide_number_label(s, 21)
print("  Slide 21: Workshop overview")

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 22 — WORKSHOP: FILE STRUCTURE
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); add_shape_bg(s, BG_DARK)
add_textbox(s, 0.8, 0.4, 11, 0.8, "Step 1: Create the File", font_size=36, bold=True, color=ACCENT3)
add_accent_line(s, 1.2, ACCENT3)

add_code_box(s, 2.5, 1.5, 8.3, 0.55, "$ touch tools_extra/threat_intel.py")

skeleton = (
    '"""Threat intelligence tools for Sysadmin Copilot.\n'
    '\n'
    'Provides IOC extraction, file hashing, and\n'
    'VirusTotal lookups. VT tools require a VT_API_KEY\n'
    'environment variable (free tier works fine).\n'
    '"""\n'
    '\n'
    'import hashlib\n'
    'import os\n'
    'import re\n'
    'from collections import Counter\n'
    '\n'
    'from langchain_core.tools import tool\n'
    'from tools import run_cmd\n'
    '\n'
    '# ... tools go here ...\n'
    '\n'
    'WRITE_TOOLS = {}  # All read-only'
)
add_code_box(s, 2.0, 2.4, 9.3, 4.5, skeleton)

slide_number_label(s, 22)
print("  Slide 22: File structure")

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 23 — WORKSHOP: IOC PATTERNS
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); add_shape_bg(s, BG_DARK)
add_textbox(s, 0.8, 0.4, 11, 0.8, "Step 2: Define IOC Patterns", font_size=36, bold=True, color=ACCENT3)
add_accent_line(s, 1.2, ACCENT3)

regex_code = (
    '_RE_IPV4   = re.compile(\n'
    '    r"\\b(?:(?:25[0-5]|2[0-4]\\d|1\\d{2}|[1-9]?\\d)\\.){3}"\n'
    '    r"(?:25[0-5]|2[0-4]\\d|1\\d{2}|[1-9]?\\d)\\b"\n'
    ')\n'
    '_RE_HASH_MD5    = re.compile(r"\\b[a-fA-F0-9]{32}\\b")\n'
    '_RE_HASH_SHA1   = re.compile(r"\\b[a-fA-F0-9]{40}\\b")\n'
    '_RE_HASH_SHA256 = re.compile(r"\\b[a-fA-F0-9]{64}\\b")\n'
    '_RE_URL    = re.compile(r"https?://[^\\s\\"\'<>]+")\n'
    '_RE_EMAIL  = re.compile(\n'
    '    r"\\b[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\\.[a-zA-Z]{2,}\\b"\n'
    ')\n'
    '_RE_DOMAIN = re.compile(\n'
    '    r"\\b(?:[a-zA-Z0-9]...)+(?:com|net|org|...)\\b"\n'
    ')\n'
    '\n'
    '# False-positive filter\n'
    '_IGNORE_IPS = {"0.0.0.0", "127.0.0.1", "255.255.255.255"}'
)
add_code_box(s, 0.8, 1.5, 8.5, 4.8, regex_code)

add_card(s, 9.8, 1.5, 3.0, 4.8)
tf = add_rich_textbox(s, 10.1, 1.8, 2.4, 4.2)
add_para(tf, "Tip:", font_size=18, bold=True, color=ACCENT3, space_after=12)
add_para(tf, "MD5 regex will match UUIDs and other 32-char hex strings.", font_size=15, color=LIGHT_GRAY, space_after=12)
add_para(tf, "Document this in the docstring!", font_size=15, bold=True, color=ACCENT, space_after=16)
add_para(tf, "We cascade: SHA256 first, exclude those from SHA1, exclude both from MD5.", font_size=15, color=LIGHT_GRAY)

slide_number_label(s, 23)
print("  Slide 23: IOC patterns")

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 24 — WORKSHOP: extract_iocs
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); add_shape_bg(s, BG_DARK)
add_textbox(s, 0.8, 0.4, 11, 0.8, "Step 3: The extract_iocs Tool", font_size=36, bold=True, color=ACCENT3)
add_accent_line(s, 1.2, ACCENT3)

ioc_code = (
    '@tool\n'
    'def extract_iocs(file_path: str,\n'
    '                 max_results: int = 50) -> str:\n'
    '    """Extract IOCs from a file.\n'
    '    Scans for IPs, domains, URLs, emails,\n'
    '    and file hashes (MD5, SHA1, SHA256).\n'
    '    """\n'
    '    # 50 MB file size guard\n'
    '    if os.path.getsize(file_path) > 50*1024*1024:\n'
    '        return "[ERROR] File too large"\n'
    '\n'
    '    text = open(file_path, "r", errors="replace").read()\n'
    '\n'
    '    ips = Counter(ip for ip in _RE_IPV4.findall(text)\n'
    '                  if ip not in _IGNORE_IPS)\n'
    '    sha256s = Counter(_RE_HASH_SHA256.findall(text))\n'
    '    sha1s = Counter(h for h in _RE_HASH_SHA1.findall(text)\n'
    '                    if h not in set(sha256s))  # cascade!\n'
    '    ...\n'
    '    # Format with counts: "185.x.x.x (x23)"'
)
add_code_box(s, 0.8, 1.5, 8.5, 5.0, ioc_code)

add_card(s, 9.8, 1.5, 3.0, 5.0)
tf = add_rich_textbox(s, 10.1, 1.8, 2.4, 4.5)
add_para(tf, "Key design:", font_size=18, bold=True, color=ACCENT, space_after=12)
add_para(tf, "@tool decorator", font_size=15, color=LIGHT_GRAY, space_after=8)
add_para(tf, "Clear docstring (the LLM reads it!)", font_size=15, color=LIGHT_GRAY, space_after=8)
add_para(tf, "50 MB file size guard", font_size=15, color=LIGHT_GRAY, space_after=8)
add_para(tf, "Counter for frequency", font_size=15, color=LIGHT_GRAY, space_after=8)
add_para(tf, "SHA256\u2192SHA1\u2192MD5 cascade", font_size=15, color=ACCENT2, bold=True, space_after=8)
add_para(tf, "Filter domains already in URLs", font_size=15, color=LIGHT_GRAY)

slide_number_label(s, 24)
print("  Slide 24: extract_iocs")

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 25 — WORKSHOP: hash_file
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); add_shape_bg(s, BG_DARK)
add_textbox(s, 0.8, 0.4, 11, 0.8, "Step 4: The hash_file Tool", font_size=36, bold=True, color=ACCENT3)
add_accent_line(s, 1.2, ACCENT3)

hash_code = (
    '@tool\n'
    'def hash_file(file_path: str) -> str:\n'
    '    """Compute MD5, SHA1, and SHA256 hashes.\n'
    '\n'
    '    Useful for verifying file integrity or\n'
    '    looking up hashes on VirusTotal.\n'
    '    """\n'
    '    md5 = hashlib.md5()\n'
    '    sha1 = hashlib.sha1()\n'
    '    sha256 = hashlib.sha256()\n'
    '\n'
    '    with open(file_path, "rb") as f:\n'
    '        for chunk in iter(\n'
    '            lambda: f.read(8192), b""\n'
    '        ):\n'
    '            md5.update(chunk)\n'
    '            sha1.update(chunk)\n'
    '            sha256.update(chunk)\n'
    '\n'
    '    return (\n'
    '        f"File:   {file_path}\\n"\n'
    '        f"MD5:    {md5.hexdigest()}\\n"\n'
    '        f"SHA1:   {sha1.hexdigest()}\\n"\n'
    '        f"SHA256: {sha256.hexdigest()}"\n'
    '    )'
)
add_code_box(s, 1.5, 1.5, 7.5, 5.2, hash_code)

add_card(s, 9.5, 2.5, 3.3, 2.5)
tf = add_rich_textbox(s, 9.8, 2.8, 2.7, 2.0)
add_para(tf, "Chunked reading", font_size=19, bold=True, color=ACCENT2, space_after=12)
add_para(tf, "8192-byte chunks \u2014 works on multi-GB files without loading into memory", font_size=16, color=LIGHT_GRAY, space_after=12)
add_para(tf, "All three hashes computed in a single pass", font_size=16, color=LIGHT_GRAY)

slide_number_label(s, 25)
print("  Slide 25: hash_file")

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 26 — WORKSHOP: VT TOOLS
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); add_shape_bg(s, BG_DARK)
add_textbox(s, 0.8, 0.4, 11, 0.8, "Step 5: VirusTotal Lookups", font_size=36, bold=True, color=ACCENT3)
add_accent_line(s, 1.2, ACCENT3)

vt_code = (
    '@tool\n'
    'def vt_hash_lookup(file_hash: str) -> str:\n'
    '    """Look up a file hash on VirusTotal\n'
    '    for malware detection results.\n'
    '\n'
    '    Requires VT_API_KEY env var (free tier ok).\n'
    '    """\n'
    '    api_key = os.environ.get("VT_API_KEY")\n'
    '    if not api_key:\n'
    '        return "[ERROR] VT_API_KEY not set."\n'
    '\n'
    '    return run_cmd([\n'
    '        "curl", "-s",\n'
    '        "-H", f"x-apikey: {api_key}",\n'
    '        f"https://www.virustotal.com/api/v3"\n'
    '        f"/files/{file_hash}",\n'
    '    ], 30)'
)
add_code_box(s, 0.8, 1.5, 7.5, 4.5, vt_code)

add_card(s, 8.8, 1.5, 4.0, 4.5)
tf = add_rich_textbox(s, 9.1, 1.8, 3.4, 4.0)
add_para(tf, "Same pattern for:", font_size=18, bold=True, color=ACCENT3, space_after=12)
add_para(tf, "vt_hash_lookup", font_size=16, color=ACCENT, font_name="Consolas", space_after=6)
add_para(tf, "\u2192 /api/v3/files/{hash}", font_size=14, color=DIM_GRAY, font_name="Consolas", space_after=10)
add_para(tf, "vt_ip_lookup", font_size=16, color=ACCENT, font_name="Consolas", space_after=6)
add_para(tf, "\u2192 /api/v3/ip_addresses/{ip}", font_size=14, color=DIM_GRAY, font_name="Consolas", space_after=10)
add_para(tf, "vt_domain_lookup", font_size=16, color=ACCENT, font_name="Consolas", space_after=6)
add_para(tf, "\u2192 /api/v3/domains/{domain}", font_size=14, color=DIM_GRAY, font_name="Consolas", space_after=14)
add_para(tf, "Free VT API tier: 4 lookups/minute", font_size=15, bold=True, color=ACCENT2)

slide_number_label(s, 26)
print("  Slide 26: VT tools")

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 27 — WORKSHOP: TESTING
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); add_shape_bg(s, BG_DARK)
add_textbox(s, 0.8, 0.4, 11, 0.8, "Step 6: Test Your Plugin", font_size=36, bold=True, color=ACCENT3)
add_accent_line(s, 1.2, ACCENT3)

test1 = (
    "\u276f tools\n"
    "\n"
    "  extract_iocs       Extract IOCs from a file\n"
    "  hash_file          Compute file hashes\n"
    "  vt_hash_lookup     Look up hash on VirusTotal\n"
    "  vt_ip_lookup       Look up IP on VirusTotal\n"
    "  vt_domain_lookup   Look up domain on VirusTotal"
)
add_code_box(s, 0.8, 1.5, 6.5, 2.8, test1)

test2 = (
    "\u276f Extract IOCs from /var/log/auth.log\n"
    "\n"
    "  [extract_iocs]\n"
    "\n"
    "  IPs (47 unique):\n"
    "    185.220.101.47  (x23)\n"
    "    192.168.1.1     (x12)\n"
    "    10.0.0.1        (x8)\n"
    "    ...\n"
    "\n"
    "  Domains (3 unique):\n"
    "    evil-scanner.ru  (x5)"
)
add_code_box(s, 0.8, 4.5, 6.5, 2.7, test2)

add_card(s, 7.8, 1.5, 4.7, 5.7)
tf = add_rich_textbox(s, 8.1, 1.8, 4.1, 5.2)
add_para(tf, "Testing checklist:", font_size=20, bold=True, color=ACCENT3, space_after=14)
add_para(tf, "1. Start the copilot", font_size=17, color=LIGHT_GRAY, space_after=8)
add_para(tf, '2. Type "tools" \u2014 new tools appear?', font_size=17, color=LIGHT_GRAY, space_after=8)
add_para(tf, "3. Ask it to extract IOCs from a log", font_size=17, color=LIGHT_GRAY, space_after=8)
add_para(tf, "4. Ask it to hash a file", font_size=17, color=LIGHT_GRAY, space_after=8)
add_para(tf, "5. If VT_API_KEY is set, try a lookup", font_size=17, color=LIGHT_GRAY, space_after=16)
add_para(tf, 'Type "audit" to verify tool calls were logged!', font_size=16, bold=True, color=ACCENT)

slide_number_label(s, 27)
print("  Slide 27: Testing")

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 28 — WORKSHOP: KEY TAKEAWAYS
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); add_shape_bg(s, BG_DARK)
add_textbox(s, 0.8, 0.4, 11, 0.8, "Plugin Recipe", font_size=38, bold=True, color=ACCENT3)
add_accent_line(s, 1.2, ACCENT3)

recipe = [
    ("1", "Create  tools_extra/your_plugin.py", ACCENT),
    ("2", "Import  @tool  from langchain_core.tools", ACCENT2),
    ("3", "Import  run_cmd  from tools", ACCENT2),
    ("4", "Write functions with clear docstrings", ACCENT),
    ("5", "Set  WRITE_TOOLS = {}  (or list write tools)", ACCENT3),
    ("6", "Restart the copilot \u2014 done!", ACCENT2),
]

y = 1.6
for num, text, clr in recipe:
    add_card(s, 1.0, y, 11.3, 0.7)
    add_circle_icon(s, 1.3, y + 0.08, 0.55, clr, num)
    add_textbox(s, 2.1, y + 0.1, 9.5, 0.5, text, font_size=20, color=LIGHT_GRAY, font_name="Calibri")
    y += 0.83

add_card(s, 2.5, y + 0.3, 8.3, 0.7, bg=RGBColor(0xE0, 0xF5, 0xE8))
add_textbox(s, 2.7, y + 0.4, 7.9, 0.5, "No core file edits.  No registration.  Just drop and go.", font_size=22, bold=True, color=ACCENT2, alignment=PP_ALIGN.CENTER)

slide_number_label(s, 28)
print("  Slide 28: Recipe")

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 29 — IDEAS
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); add_shape_bg(s, BG_DARK)
add_textbox(s, 0.8, 0.4, 11, 0.8, "Ideas for More Plugins", font_size=38, bold=True, color=TITLE_TEXT)
add_accent_line(s, 1.2)

ideas = [
    ("Docker container management", ACCENT),
    ("Kubernetes pod inspection", ACCENT2),
    ("SSL certificate checker", ACCENT3),
    ("Backup status monitor", ACCENT),
    ("Cloud provider tools (AWS / GCP status)", ACCENT2),
    ("Custom application health checks", ACCENT3),
    ("Database query tools", ACCENT),
]

# Two-column layout
col1 = ideas[:4]
col2 = ideas[4:]

y = 1.8
for text, clr in col1:
    add_card(s, 0.8, y, 5.7, 0.65)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(y), Inches(0.12), Inches(0.65))
    bar.fill.solid(); bar.fill.fore_color.rgb = clr; bar.line.fill.background()
    add_textbox(s, 1.2, y + 0.1, 5, 0.4, text, font_size=18, color=LIGHT_GRAY)
    y += 0.82

y = 1.8
for text, clr in col2:
    add_card(s, 6.8, y, 5.7, 0.65)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(y), Inches(0.12), Inches(0.65))
    bar.fill.solid(); bar.fill.fore_color.rgb = clr; bar.line.fill.background()
    add_textbox(s, 7.2, y + 0.1, 5, 0.4, text, font_size=18, color=LIGHT_GRAY)
    y += 0.82

add_textbox(s, 0.8, 5.5, 11.5, 0.8, "The plugin system makes it easy to adapt the copilot\nto your specific environment and workflows.", font_size=18, color=DIM_GRAY, alignment=PP_ALIGN.CENTER)

slide_number_label(s, 29)
print("  Slide 29: Ideas")

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 30 — THANK YOU
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); add_shape_bg(s, BG_DARK)

# Top accent bar
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

add_textbox(s, 1, 2.0, 11.3, 1.2, "Thank You!", font_size=56, bold=True, color=TITLE_TEXT, alignment=PP_ALIGN.CENTER)
add_textbox(s, 1, 3.3, 11.3, 0.8, "Questions?", font_size=32, color=ACCENT, alignment=PP_ALIGN.CENTER)

add_accent_line(s, 4.3)

add_textbox(s, 1, 4.8, 11.3, 0.5, "All code shown today is in the repo", font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_code_box(s, 3.5, 5.5, 6.3, 0.55, "$ sysadmin-copilot   # try it yourself!")

print("  Slide 30: Thank you")

# ══════════════════════════════════════════════════════════════════════════
#  SAVE
# ══════════════════════════════════════════════════════════════════════════
output_path = "/home/clint/Documents/Projects/AI/Sysadmin-Copilot/presentation.pptx"
prs.save(output_path)
print(f"\nSaved: {output_path}")
print(f"Total slides: {len(prs.slides)}")

cleanup_images()
